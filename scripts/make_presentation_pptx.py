"""make_presentation_pptx.py — the deliverable deck, openable in PowerPoint.
EXTREMELY detailed edition: every functionality area gets its own slide,
following PRESENTATION_HANDOFF.md, with the shot-pack images embedded.
Author: Dimitres Kisimov.

    python scripts/make_presentation_pptx.py
Out: deliverable/local_only/SCS_Studio_Presentation_v3.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parents[1]
SHOTS = REPO / "deliverable" / "local_only" / "presentation_shots"
OUT = REPO / "deliverable" / "local_only" / "SCS_Studio_Presentation_v3.pptx"

INK = RGBColor(0x1F, 0x27, 0x33)
BRAND = RGBColor(0x2F, 0x6B, 0xFF)
MUT = RGBColor(0x6B, 0x76, 0x88)


def add_slide(prs, title, bullets=None, image=None, note=None, kicker=None):
    """image: filename, or a list of up to 2 filenames stacked vertically."""
    s = prs.slides.add_slide(prs.slide_layouts[6])       # blank
    if kicker:
        kx = s.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.3))
        kp = kx.text_frame.paragraphs[0]
        kp.text = kicker.upper()
        kp.font.size = Pt(12)
        kp.font.bold = True
        kp.font.color.rgb = BRAND
    tx = s.shapes.add_textbox(Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.9))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK
    imgs = [image] if isinstance(image, str) else (image or [])
    imgs = [i for i in imgs if (SHOTS / i).exists()]
    body_w = Inches(6.4) if imgs else Inches(12.3)
    if bullets:
        bx = s.shapes.add_textbox(Inches(0.5), Inches(1.35), body_w, Inches(5.7))
        tf = bx.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            cont = b.startswith(("—", " ", "→"))
            para.text = b if cont else ("• " + b)
            para.font.size = Pt(13)
            para.font.color.rgb = MUT if cont else INK
            para.space_after = Pt(5)
    if imgs:
        h = Inches(5.5) if len(imgs) == 1 else Inches(2.7)
        for i, img in enumerate(imgs[:2]):
            s.shapes.add_picture(str(SHOTS / img), Inches(7.15),
                                 Inches(1.35) + i * (h + Inches(0.15)),
                                 width=Inches(5.7))
    if note:
        nb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
        np_ = nb.text_frame.paragraphs[0]
        np_.text = note
        np_.font.size = Pt(10)
        np_.font.color.rgb = MUT
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- 1 · Title -------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tx = s.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(11), Inches(2.6))
    tf = tx.text_frame
    tf.paragraphs[0].text = "SCS Studio"
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BRAND
    p2 = tf.add_paragraph()
    p2.text = "photo → 3D → room → BIM   ·   one photo in, a compliant furnished building out"
    p2.font.size = Pt(22)
    p2.font.color.rgb = INK
    p3 = tf.add_paragraph()
    p3.text = "Dimitres Kisimov · v3.0 · 2026 · github.com/Dimitres-Kisimov/3DpicToIFCModeling"
    p3.font.size = Pt(15)
    p3.font.color.rgb = MUT

    # ---- 2 · Executive summary ------------------------------------------
    add_slide(prs, "What SCS Studio is — in one breath", [
        "A local web app (Node + Python, localhost:3000) started by double-clicking SCS_Studio.bat",
        "A furniture PHOTO becomes a repaired, real-dimensioned 3D catalog item",
        "A constraint solver (Google CP-SAT) furnishes single rooms AND whole IFC buildings",
        "Placement obeys German ASR workplace law + DIN 277 room usage + human logic",
        "Results export back to BIM: IFC for Revit, GLB for 3D, CSV schedules",
        "Every claim is test-gated; a built-in Research hub documents the evidence",
        "Shipped: v1.0.0 → v2.0.1 (public GitHub release, 512 MB) → v3.0.0"],
        image="20_research_hub_7_steps.jpg", kicker="Executive summary")

    # ---- 3 · Problem -----------------------------------------------------
    add_slide(prs, "The problem: furnishing BIM by hand", [
        "Every object is placed one by one: search it, drag it in, rotate it,",
        "  check the clearances, repeat — thousands of decisions per office building",
        "German workplace law (ASR) adds LEGAL constraints most tools ignore:",
        "  workstation areas, movement zones, escape-route widths, safety access",
        "Checking compliance is manual eyeballing — and must be redone after every edit",
        "Result: days to weeks of skilled manual work per building, error-prone"],
        image="12_tab3_building.jpg", kicker="The problem")

    # ---- 4 · Time savings ------------------------------------------------
    add_slide(prs, "Time savings — measured, not estimated", [
        "Furnish ONE office room legally (10–20 pieces):",
        "  by hand 15–45 min  →  SECONDS, one click",
        "Furnish a 6-storey office — 2,191 pieces (Autodesk 210 King, Toronto):",
        "  by hand: days  →  ~17 min first run, ~12 min warm",
        "Furnish the whole 15-building fleet (~6,000 pieces):",
        "  by hand: weeks  →  ~16 minutes total",
        "Re-check every clearance after any edit:",
        "  manual eyeballing  →  automatic and instant, on every move"],
        note="Measured on the real fleet, warm caches; a first-ever scan of a huge IFC adds a one-time cost.",
        kicker="The core pitch")

    # ---- 5 · Pipeline ----------------------------------------------------
    add_slide(prs, "The pipeline in one picture", [
        "Tab 1 — Generate object: photo → AI mesh → repair → catalog item",
        "Tab 2 — Build a room: the engine hands-on; the quick RUN TEST",
        "Tab 3 — Building: the BULK WORK — same logic across a whole IFC",
        "Research hub: 8 guided steps of evidence behind every claim",
        "One shared 3D viewport (xeokit + model-viewer); 2D plan editors synced live",
        "Everything runs locally — no cloud dependency at runtime"],
        image="30_roadmap_dated.jpg", kicker="Architecture")

    # ---- 6 · Mechanism 1: knows what belongs where -----------------------
    add_slide(prs, "Mechanism 1 — the algorithm knows what belongs where", [
        "A rulebook maps every room type to its permissible furniture:",
        "  offices ≠ kitchens ≠ meeting rooms ≠ server rooms ≠ balconies",
        "DIN 277 usage groups decide what gets furnished at all —",
        "  technical, circulation and sanitary rooms are correctly LEFT EMPTY",
        "German ASR workplace law is cited from the legal text (rule_packs.py)",
        "Human logic on top: chairs face tables, the monitor faces the sitter,",
        "  the microwave sits ON the counter, the fridge belongs in the kitchen",
        "What cannot sit legally is REFUSED and reported — never forced"],
        image="23_item_register_rules.jpg", kicker="Mechanism 1 / 5")

    # ---- 7 · Human placement logic detail --------------------------------
    add_slide(prs, "Human placement logic — the details that sell it", [
        "Waste bin: at arm's reach BESIDE the desk (~0.2 m) — it is the one",
        "  'easily movable' object, so it gets minimal clearance, not a fortress",
        "Planters: corners and window zones only — never mid-room",
        "Projector: CEILING-mounted, aimed at its screen, ASR-safe headroom,",
        "  throw distance scaled to the screen size",
        "Presentation rows: chairs FACE the whiteboard/screen, centered on its",
        "  axis, at the ASR-required viewing distance — including partial last rows",
        "Armchairs: in ROWS with a side table — 2 face the same way;",
        "  4 = two opposed pairs; never blocking the walkway",
        "Rectangular tables: edge seating, max 4 per side placed centrally;",
        "  spare chairs park against the wall facing the room",
        "Round tables only get the 'petal' ring — rectangles never do"],
        image=["35_3d_meeting_hub.jpg", "36_3d_break_social.jpg"],
        kicker="Mechanism 1 · continued")

    # ---- 8 · Safety ------------------------------------------------------
    add_slide(prs, "Safety is non-negotiable — even at extreme density", [
        "Fire extinguisher and first-aid cabinet get PROTECTED access strips:",
        "  a keep-out zone in front that no furniture may ever enter",
        "Enforced at every density — Dense cannot bury the extinguisher",
        "Whiteboards / presentation screens get visibility strips the same way",
        "Wall-mounted safety gear seeks a free wall slot automatically",
        "  (extinguisher 1.00 m, first-aid 1.35 m, ASR mounting heights)",
        "0.90 m person-path to EVERY placed item is solver-verified",
        "Escape-route widths by occupancy: 0.90 / 1.00 / 1.20 m (ASR A2.3/A1.8)"],
        image="33_2d_office_densest_plan.jpg", kicker="Mechanism 1 · safety")

    # ---- 9 · Room census / bilingual -------------------------------------
    add_slide(prs, "1,506 rooms auto-classified — in four languages", [
        "The fleet's 1,506 rooms were typed automatically from their names",
        "German, English, Dutch and French all resolve to the same logic:",
        "  Büro = Tenant = office · Kellerraum = storage · Garderobe = wardrobe",
        "Grounded in DIN 277: NUF (usable) furnished, TF/VF (technical/",
        "  circulation) + sanitary correctly skipped",
        "896 rooms auto-furnished · 562 correctly left empty · 48 garbage labels",
        "  (manual picks cover those)",
        "6 dedicated room types added: storage, wardrobe, entry/foyer, balcony,",
        "  print/copy, IT/server — each with its own furniture pack",
        "Unknown room? Clean-slate mode: the user picks a type or chooses manually"],
        image="26_human_layouts_densities.jpg", kicker="Mechanism 1 · scale")

    # ---- 10 · Templates --------------------------------------------------
    add_slide(prs, "Mechanism 2 — templates: furnish once, apply everywhere", [
        "Right-click a room card → its furniture picks become the DEFAULT",
        "  for every room of that name/type in the building",
        "📐 Layout blueprint: copy a finished room's EXACT arrangement",
        "  to all similar rooms in one action",
        "Identical-size rooms match 1:1 — same layout, instantly",
        "Slightly different sizes are adapted piece by piece WITH clash-checking",
        "Too-small rooms are refused with a clear message — no silent breakage",
        "Real effect: furnish one hotel room → get the whole corridor"],
        kicker="Mechanism 2 / 5")

    # ---- 11 · Multi-copy -------------------------------------------------
    add_slide(prs, "Mechanism 3 — multi-copy: counts, not clicks", [
        "Right-click a furniture chip and type a count",
        "10 chairs in ONE action instead of ten separate selections",
        "The solver places all copies legally in the same run",
        "Combines with templates: a counted set becomes the type default too"],
        kicker="Mechanism 3 / 5")

    # ---- 12 · Densities --------------------------------------------------
    add_slide(prs, "Mechanism 4 — one click, three densities", [
        "Light / Medium / Dense population for a room or a WHOLE building",
        "Dense staffs toward the ASR legal maximum — never past it:",
        "  §5(3) hard cap: 8 m² first workstation + 6 m² each further",
        "Capacity guard: suggested sets are trimmed to what physically and",
        "  legally fits (usable area × footprint × 1.5 walking factor)",
        "Density is monotonic — verified by a sweep test: Light ⊆ Medium ⊆ Dense",
        "Suggestions come from the SAME engine buildings use — the room tab's",
        "  ✨Suggest is a truthful preview of bulk behavior"],
        image="24_live_layout_visualizer.jpg", kicker="Mechanism 4 / 5")

    # ---- 13 · Manual control ---------------------------------------------
    add_slide(prs, "Mechanism 5 — full manual control on top of automation", [
        "MOVE: drag any piece in 3D or in the 2D plan — legality is checked",
        "  live; illegal spots flag red and snap back",
        "ROTATE: R key / ⟳ button / Ctrl+right-click — 90° yaw steps, so an",
        "  object faces any of the four walls; ceiling/floor tilts are impossible",
        "DELETE: Delete key / 🗑 button — the item vanishes from 3D, 2D",
        "  AND every export (GLB, IFC, CSV)",
        "LOCK/UNLOCK: protect a finished plan from accidental drags —",
        "  in the building editor and the room editor alike",
        "Every manual change is re-validated and BAKED into the exports:",
        "  the IFC you download holds final positions and rotations"],
        image=["31_2d_office_team_plan.jpg", "32_3d_office_team_render.jpg"],
        kicker="Mechanism 5 / 5")

    # ---- 14 · Tab 1 ------------------------------------------------------
    add_slide(prs, "Tab 1 — Generate object: photo → catalog item", [
        "Drag & drop ONE furniture photo (or 'Try a sample chair')",
        "Quality modes: Detailed = TripoSR builds a new mesh locally (~60–90 s);",
        "  Fast = classify + best catalog match instantly",
        "Declared object type forces category + IFC class + numbering",
        "Chair-base graft: AI can't rebuild thin legs from one photo — we graft",
        "  clean parametric bases (4 styles) onto the AI-generated seat",
        "Auto-registered with a professional code: item – engine – number",
        "  (desk-SAM3D-001); deleting renumbers the catalog automatically",
        "Real dimensions assigned per category — a desk lands desk-sized"],
        image="10_tab1_generate_object.jpg", kicker="Tab 1")

    # ---- 15 · Engines ----------------------------------------------------
    add_slide(prs, "The AI engines — benchmarked honestly", [
        "Five engines: TripoSR (runs LOCALLY on the app machine),",
        "  TripoSG, TRELLIS, SAM 3D, InstantMesh (cloud GPUs)",
        "187 internet photos, identical inputs, F-score vs ground-truth meshes",
        "TripoSG best at 0.393; full A/B lists browsable in 3D in the hub",
        "License-screened: EU-usable engines only (documented decision)",
        "Every engine's output is inspectable side by side — no cherry-picking"],
        image=["21_benchmark_ab_lists.jpg", "22_multi_ai_visualizer.jpg"],
        kicker="Research")

    # ---- 16 · Repair -----------------------------------------------------
    add_slide(prs, "Mesh repair — broken AI output becomes usable furniture", [
        "Raw AI meshes arrive holey, noisy and non-manifold",
        "Repair pipeline: hole filling, smoothing, decimation, re-orientation",
        "Parametric chair-base grafts replace unreconstructable thin geometry",
        "Output is watertight, real-dimensioned, IFC-classifiable",
        "Before/after comparisons in the research hub for every repair stage"],
        kicker="Research")

    # ---- 17 · Catalog ----------------------------------------------------
    add_slide(prs, "The catalog — 38 categories, hundreds of variants", [
        "Amazon Berkeley Objects (CC-BY-4.0): 4.7 GB of real product meshes",
        "  with photos — chairs, desks, sofas, shelving, lamps…",
        "Our own AI-generated items live beside them, same numbering system",
        "44 extra screened variants for 15 thin categories (strict rule:",
        "  a 'coffee machine' variant IS a coffee machine, never a lookalike)",
        "CC0 assets + parametric styles fill the gaps",
        "Custom categories: upload your OWN IFC furniture and it becomes",
        "  a first-class, selectable catalog category",
        "Variant browser with product photos on every category"],
        image="25_showroom_all_items.jpg", kicker="Catalog")

    # ---- 18 · Tab 2 ------------------------------------------------------
    add_slide(prs, "Tab 2 — Build a room: the engine, hands-on", [
        "Room type (office-first, ASR-tagged) × size presets × density",
        "✨Suggest asks the SAME engine buildings use — capacity-guarded",
        "  so the suggestion always fits the room legally",
        "Pick items per category from the variant browser, set counts",
        "CP-SAT solver places everything under ASR + human logic in seconds",
        "Positioning: this tab is the quick RUN TEST / demo of the engine —",
        "  the bulk production value is the Building tab",
        "Exports: CSV schedule · GLB scene · BIM IFC"],
        image="11_tab2_build_a_room.jpg", kicker="Tab 2")

    # ---- 19 · Collision checker ------------------------------------------
    add_slide(prs, "Collision checker with walking paths (room tab)", [
        "Every manual move re-validates TWO things live:",
        "  1. overlaps — no two pieces may intersect",
        "  2. the 0.90 m person-path — every item must stay REACHABLE",
        "Violations flag red and name the unreachable item",
        "No items outside the room: the exact polygon meter is a test gate",
        "Live-verified for this release: shoving desk-2 onto desk-1 flagged",
        "  violations=True with the circulation path in the response"],
        image="31_2d_office_dense_plan.jpg", kicker="Tab 2 · v3")

    # ---- 20 · Room 2D editor ---------------------------------------------
    add_slide(prs, "The room 2D editor — same mechanics as the building", [
        "Drag any piece in the plan — the 3D scene follows",
        "Rotate: R key / ↻ — 90° steps toward the four walls",
        "DELETE (v3): Delete key removes the item from the schedule, the",
        "  3D scene, the zones AND the rebuilt IFC — all views update",
        "LOCK/UNLOCK (v3): one button freezes dragging on a finished plan",
        "Banner confirms: 'removed from the room — 3D, IFC and exports updated'"],
        image="37_2d_showroom_all_items_plan.jpg", kicker="Tab 2 · v3")

    # ---- 21 · Tab 3 upload -----------------------------------------------
    add_slide(prs, "Tab 3 — Building: drop ANY .ifc", [
        "Instant profile on upload: storeys, spaces, products",
        "Automatic floor dissection — no manual setup, no naming conventions",
        "Works on files the app has NEVER seen (drop-tested, next slides)",
        "15 licensed buildings ship as the demo fleet: homes, university",
        "  institutes, corporate offices, synthetic towers",
        "Uploaded buildings get the FULL catalog, including your own",
        "  generated items — verified end-to-end"],
        image="12_tab3_building.jpg", kicker="Tab 3")

    # ---- 22 · Floor-by-floor ---------------------------------------------
    add_slide(prs, "Floor-by-floor workflow", [
        "Floor chips isolate each storey in the 3D view AND the 2D plan",
        "Double-click a room in 2D → the camera flies into it in 3D",
        "Per-room cards show type, area and smart suggestions",
        "Templates + blueprints + multi-copy work per floor or building-wide",
        "Clean-slate manual picks for any room that should differ",
        "X-ray mode ghosts the shell — inspect furniture through the walls",
        "Lock camera for stable presentation views"],
        image="27_xray_fleet.jpg", kicker="Tab 3")

    # ---- 23 · Building manual edit ----------------------------------------
    add_slide(prs, "Editing a populated building — piece by piece", [
        "Select any placed piece in 3D: drag it (legality live), rotate it",
        "  90° (R / ⟳ / Ctrl+right-click), or Delete it",
        "The 2D floor plan mirrors every change instantly — and vice versa",
        "Deleting overcrowded objects: the capacity guard already trims at",
        "  Dense, and anything else is one keypress away",
        "Colored, real-dimensioned pieces — chairs look like chairs",
        "All edits survive into the exports at their final positions"],
        image=["33_2d_office_open_plan.jpg", "34_3d_office_open_render.jpg"],
        kicker="Tab 3")

    # ---- 24 · Exports ----------------------------------------------------
    add_slide(prs, "Exports — back to BIM, faithfully", [
        "💾 Whole-building GLB for any 3D pipeline",
        "📥 BIM IFC: every piece is an IfcFurnishingElement in its correct",
        "  storey, at its FINAL position — moves, 90° rotations and deletions",
        "  are baked in (rotation as a proper yaw matrix, deleted pieces gone)",
        "Room tab: CSV schedule + GLB + IFC per room",
        "Round-trip: the exported IFC re-opens in Revit and in the app itself"],
        kicker="Exports")

    # ---- 25 · ASR --------------------------------------------------------
    add_slide(prs, "German ASR compliance — cited, not claimed", [
        "ArbStättV §5(3): 8 m² for the first workstation + 6 m² each further —",
        "  enforced as a HARD staffing cap per room",
        "ASR A1.2: ≥1.5 m² free movement area, ≥1.00 m depth at every desk",
        "ASR A1.8/A2.3: route widths 0.90 / 1.00 / 1.20 m by occupancy;",
        "  we enforce a 0.90 m access path to EVERY item, stricter than the",
        "  0.60 m workstation-access minimum",
        "Safety equipment freely accessible — protected keep-out strips",
        "Three enforcement tiers by room type (presentation/reception 1.20 m,",
        "  break 1.00 m, quiet 0.90 m)",
        "Full mapping in docs/ASR_COMPLIANCE.md: every rule is ASR-cited,",
        "  STRICTER than ASR, or marked as practice"],
        image="23_item_register_rules.jpg", kicker="Compliance")

    # ---- 26 · Item Logic Register ----------------------------------------
    add_slide(prs, "The Item Logic Register — every rule, auditable", [
        "One table: every furniture item × its placement logic × its",
        "  relationships × the German-standard verdict",
        "Who sits beside whom: chair→table, bin→desk, projector→screen,",
        "  planter→corner/window, extinguisher→wall+strip…",
        "Rendered three ways: in-app page, Markdown, Word (.docx)",
        "This is the contract the solver is tested against"],
        image="23_item_register_rules.jpg", kicker="Compliance")

    # ---- 27 · Drop test --------------------------------------------------
    add_slide(prs, "Drop test: a 7-storey IFC the app never saw", [
        "A fresh 12.9 MB tower, through the real upload path",
        "Instant profile: 10 storeys, 81 spaces, 720 products — all correct",
        "Room dissection: 9/9 space-bearing floors, 81/81 rooms EXACT",
        "  against the IFC ground truth",
        "Full catalog usable per room, including our generated items",
        "Repeated three times, documented in docs/DROP_TEST_REPORT.md",
        "Conclusion: ANY reasonable IFC drops in and works"],
        image="50_building_04_SCS_Tower_7_storeys_synthetic.jpg",
        kicker="Evidence")

    # ---- 28 · Fleet ------------------------------------------------------
    add_slide(prs, "The fleet at scale — 15 buildings, ~6,000 pieces", [
        "Homes, university institutes, corporate offices, synthetic towers —",
        "  incl. Autodesk's 210 King Toronto (6 storeys, 2,191 pieces)",
        "~2,237 machine-checked human connections exported as data:",
        "  chair→desk, projector→screen, bin-beside-desk, door flanks…",
        "X-Ray Fleet page: every building's furniture through ghosted walls",
        "Fleet Report: per-building tables of items, connections, compile time",
        "Exact-geometry clash meter over the whole fleet: ZERO real overlaps",
        "  (two 0.02–0.05 m² slivers of the movable-bin class, disclosed)"],
        image=["29_buildings_showcase.jpg", "28_fleet_report_connections.jpg"],
        kicker="Evidence")

    # ---- 29 · Building gallery -------------------------------------------
    add_slide(prs, "Every building, populated and X-rayed", [
        "The shot pack holds an X-ray of ALL 15 buildings (50_building_00–14)",
        "From a Dutch duplex to an 8-floor institute tower — same engine,",
        "  same rules, no per-building tuning",
        "210 King: days of manual work → ~17 minutes, 2,191 legal pieces",
        "Use these as full-bleed image slides in the final deck"],
        image=["50_building_00_210_King_Autodesk_Toronto_office_6.jpg",
               "50_building_07_HHS_Office_institute_labs_workshop.jpg"],
        kicker="Evidence")

    # ---- 30 · Test gates -------------------------------------------------
    add_slide(prs, "Test evidence — the gates this release passed", [
        "End-to-end smoke: 8/8 user capabilities, repeated runs,",
        "  including a LIVE photo→3D generation each time",
        "Floor dissection: 15/15 buildings exact vs IFC ground truth",
        "Drop test: 9/9 floors, 81/81 rooms on a never-seen tower",
        "Ergonomics gate: 40/40 checks × 2 runs — row facing, mounting",
        "  heights, microwave ON counter, door flanks, 0 clashes",
        "Density monotonicity sweep: Light ⊆ Medium ⊆ Dense",
        "Exact-geometry meter over ~6,000 pieces: zero real overlaps",
        "Release bundle boot-verified from a clean extraction",
        "Room census: 1,506 rooms — 896 furnished, 562 correctly empty"],
        kicker="Evidence")

    # ---- 31 · Research hub ------------------------------------------------
    add_slide(prs, "The Research hub — evidence built into the product", [
        "8 guided steps with arrows — a narrated tour of the whole system",
        "Step 1: 'Code files & what they do' — EVERY button, right-click and",
        "  shortcut mapped to the exact source file that implements it",
        "Then: benchmarks → repair → catalog → room logic (Item Register) →",
        "  buildings (X-Ray Fleet, Fleet Report) → proof (reports, roadmap)",
        "A reviewer can verify any slide of this deck inside the app"],
        image="20_research_hub_7_steps.jpg", kicker="Transparency")

    # ---- 32 · How to get it ----------------------------------------------
    add_slide(prs, "How to get it — install in two clicks", [
        "GitHub (public): github.com/Dimitres-Kisimov/3DpicToIFCModeling/releases",
        "Download the *_lite.zip (~512 MB) → unzip → double-click",
        "  SCS_Studio.bat — it self-installs Node/Python deps on first run",
        "The full furniture catalog (4.7 GB, ABO CC-BY-4.0) ships separately:",
        "  TransferXL link → extract INTO the app folder so files land at",
        "  SCS_Studio/data/mesh_library_abo/ → restart the app",
        "After that the machine is IDENTICAL to the development machine",
        "No-TransferXL alternative: python backend/python-scripts/",
        "  download_abo_subset.py rebuilds the library from Amazon's public data"],
        image="40_github_release_v1_tag.jpg", kicker="Distribution")

    # ---- 33 · Future vision ----------------------------------------------
    add_slide(prs, "Future vision — training AI on these layouts", [
        "The fleet report already exports every placement RELATION as",
        "  machine-checked data — 2,237 labeled examples and growing",
        "That is a training set: imitation learning over the solver's decisions",
        "Next step: generative simulation of ENTIRE buildings —",
        "  propose the usage mix per floor, simulate occupancy flows,",
        "  iterate layouts against ASR automatically",
        "The pipeline is already structured for it:",
        "  rules → solver → machine-readable outcomes → (next) learned policy"],
        kicker="Vision")

    # ---- 34 · Version history --------------------------------------------
    add_slide(prs, "Release history", [
        "v1.0.0 — the full pipeline: photo→3D, rooms, buildings, exports",
        "v2.0.1 — public release: 90° rotation everywhere, delete-everywhere,",
        "  safety access strips, armchair rows, rect-table edge seating,",
        "  capacity-guarded suggestions, kitchen logic, 512 MB lite bundle",
        "v2.1.0 — DIN 277 bilingual lexicon: 1,506 rooms, 6 new room types",
        "v3.0.0 — room-tab parity (delete + lock/unlock + collision checker),",
        "  presentation handoff, this deck, the catalog transfer package"],
        kicker="Versions")

    # ---- 35 · Thanks -----------------------------------------------------
    add_slide(prs, "Thank you", [
        "SCS Studio v3.0 — photo → 3D → room → BIM",
        "Author: Dimitres Kisimov",
        "Data: Amazon Berkeley Objects (CC-BY-4.0), PolyHaven (CC0)",
        "Buildings: Autodesk / buildingSMART / TU Eindhoven sample IFCs",
        "  + synthetic towers (all licenses documented in the repo)",
        "All evidence live in the app's Research hub — ask for a demo"],
        kicker="Credits")

    prs.save(str(OUT))
    print(f"deck: {len(prs.slides._sldIdLst)} slides -> {OUT}")


if __name__ == "__main__":
    main()
