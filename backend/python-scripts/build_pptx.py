"""
build_pptx.py — bilingual (German | English) two-part presentation for 3DpicToIFCModeling.

Part 1 (Gülriz, analytical): intro, objective/scope, single-view limitations, energy/cost,
AI approaches, HuggingFace + licences, the pivot.
Part 2 (Dimitrius, technical): system overview, catalog, layout engine, results figures,
capacity, web app, accuracy evaluation, bake-off, conclusion.

Every content slide is split-screen: Deutsch (left) | English (right). Figures reused from
paper_figures/. Run: python backend/python-scripts/build_pptx.py
Out: deliverable/SCS_Presentation_2026-06-23.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

REPO = Path(__file__).resolve().parents[2]
FIG = REPO / "paper_figures"
OUT = REPO / "deliverable" / "SCS_Presentation_2026-06-23.pptx"

ACCENT = RGBColor(0x2F, 0x81, 0xF7)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEF, 0xF3, 0xFA)

EMU_IN = 914400
SW, SH = 13.333, 7.5


def _txt(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def _line(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT, space=4):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    indent = 0
    while text.startswith("  "):
        indent += 1; text = text[2:]
    p.level = min(indent, 4)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p


def _bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


def _accent_bar(slide, y=0.0, h=0.12):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(y), Inches(SW), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()


def _header(slide, title_de, title_en):
    _accent_bar(slide, 0, 0.14)
    tf = _txt(slide, 0.5, 0.28, SW - 1, 1.0)
    _line(tf, title_de, 25, DARK, bold=True, first=True, space=0)
    _line(tf, title_en, 17, ACCENT, bold=False, space=0)


def _footer(slide, presenter, num):
    tf = _txt(slide, 0.5, SH - 0.42, SW - 1, 0.3)
    p = _line(tf, f"{presenter}", 9, MUTED, first=True)
    tf2 = _txt(slide, SW - 1.2, SH - 0.42, 0.9, 0.3)
    _line(tf2, str(num), 9, MUTED, first=True, align=PP_ALIGN.RIGHT)


def _flag(slide, x, y, label):
    tf = _txt(slide, x, y, 3.0, 0.3)
    _line(tf, label, 11, ACCENT, bold=True, first=True, space=2)


def _column(slide, x, y, w, h, label, lines, size=14):
    _flag(slide, x, y, label)
    tf = _txt(slide, x, y + 0.34, w, h)
    for i, ln in enumerate(lines):
        bullet = ln if ln.startswith("  ") else ("• " + ln if ln else "")
        _line(tf, bullet, size, DARK, first=(i == 0), space=5)


def _fit(path, box_w, box_h):
    im = Image.open(path); iw, ih = im.size
    r = min(box_w / iw, box_h / ih)
    return iw * r, ih * r


def _image(slide, path, cx, cy, box_w, box_h):
    w, h = _fit(path, box_w, box_h)
    slide.shapes.add_picture(str(path), Inches(cx - w / 2), Inches(cy - h / 2),
                             Inches(w), Inches(h))


# ---------------------------------------------------------------- slide builders
def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _accent_bar(s, 0, 0.18); _accent_bar(s, SH - 0.18, 0.18)
    tf = _txt(s, 0.8, 2.1, SW - 1.6, 2.0, anchor=MSO_ANCHOR.TOP)
    _line(tf, "Vom Foto zum BIM-Modell", 40, DARK, bold=True, first=True, align=PP_ALIGN.CENTER, space=2)
    _line(tf, "From Photo to BIM Model", 26, ACCENT, align=PP_ALIGN.CENTER, space=10)
    tf2 = _txt(s, 0.8, 4.2, SW - 1.6, 1.6)
    _line(tf2, "KI-gestützte Raummöblierung & Foto→3D-Pipeline für IFC/BIM", 15, MUTED,
          first=True, align=PP_ALIGN.CENTER, space=2)
    _line(tf2, "AI room population & photo→3D pipeline for IFC/BIM", 15, MUTED, align=PP_ALIGN.CENTER, space=14)
    _line(tf2, "Gülriz  ·  Dimitrius      |      23. Juni 2026", 14, DARK, bold=True,
          align=PP_ALIGN.CENTER, space=0)


def slide_divider(prs, part_de, part_en, presenter):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, RGBColor(0x14, 0x18, 0x20))
    _accent_bar(s, 3.2, 0.06)
    tf = _txt(s, 0.8, 2.5, SW - 1.6, 2.2, anchor=MSO_ANCHOR.MIDDLE)
    _line(tf, part_de, 30, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, space=4)
    _line(tf, part_en, 20, RGBColor(0x8F, 0xB8, 0xF7), align=PP_ALIGN.CENTER, space=12)
    _line(tf, presenter, 16, RGBColor(0xC8, 0xD0, 0xDA), align=PP_ALIGN.CENTER, space=0)


def slide_content(prs, t_de, t_en, de, en, presenter, num, img=None, cap_de="", cap_en="", size=14):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, t_de, t_en)
    col_h = 3.0 if img else 5.1
    _column(s, 0.5, 1.5, 5.9, col_h, "Deutsch", de, size)
    _column(s, 6.9, 1.5, 5.9, col_h, "English", en, size)
    # divider line between columns
    ln = s.shapes.add_shape(1, Inches(6.65), Inches(1.55), Pt(1.2), Inches(col_h + 0.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = LIGHT; ln.line.fill.background()
    if img:
        _image(s, FIG / img, SW / 2, 5.55, SW - 3.0, 2.4)
        if cap_de or cap_en:
            tf = _txt(s, 0.5, 6.85, SW - 1, 0.5)
            _line(tf, cap_de, 10, MUTED, first=True, align=PP_ALIGN.CENTER, space=0)
            _line(tf, cap_en, 10, ACCENT, align=PP_ALIGN.CENTER, space=0)
    _footer(s, presenter, num)


def slide_figure(prs, t_de, t_en, img, cap_de, cap_en, presenter, num, de=None, en=None, img2=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, t_de, t_en)
    if img2:
        _image(s, FIG / img, SW * 0.27, 3.7, SW * 0.46, 3.6)
        _image(s, FIG / img2, SW * 0.73, 3.7, SW * 0.46, 3.6)
    else:
        _image(s, FIG / img, SW / 2, 3.9, SW - 2.2, 4.4)
    yb = 6.55
    if de or en:
        _column(s, 0.5, 1.45, 5.9, 0.9, "Deutsch", de or [], 12)
        _column(s, 6.9, 1.45, 5.9, 0.9, "English", en or [], 12)
    tf = _txt(s, 0.5, yb, SW - 1, 0.6)
    _line(tf, cap_de, 12, MUTED, bold=True, first=True, align=PP_ALIGN.CENTER, space=1)
    _line(tf, cap_en, 12, ACCENT, align=PP_ALIGN.CENTER, space=0)
    _footer(s, presenter, num)


# ---------------------------------------------------------------- build
def main():
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN)); prs.slide_height = Emu(int(SH * EMU_IN))
    G, D = "Gülriz", "Dimitrius"
    n = [1]
    def nxt(): n[0] += 1; return n[0]

    slide_title(prs)

    # ---- PART 1 — Gülriz ----
    slide_divider(prs, "Teil 1 — Analyse & Ansatz", "Part 1 — Analysis & Approach", "Gülriz")

    slide_content(prs, "Einführung", "Introduction",
        ["Foto → Objekttabelle → KI-Raumlayout → IFC/BIM + 3D-Betrachter",
         "Jedes fotografierte Objekt wird eine Zeile: Typ, Maße, Material, 3D-Modell, Lizenz",
         "Die Objekttabelle ist die zentrale Wahrheitsquelle",
         "Export als CSV / IFC; eine KI möbliert den Raum funktional"],
        ["Photo → object table → AI room layout → IFC/BIM + 3D viewer",
         "Each photographed object becomes a row: type, size, material, 3D model, licence",
         "The object table is the single source of truth",
         "Export as CSV / IFC; an AI furnishes the room functionally"], G, nxt())

    slide_content(prs, "Ziel & Umfang", "Objective & Scope",
        ["Ziel: Bild→IFC-Pipeline für Büromöbel (kein ganzes Gebäude-BIM)",
         "Umfang: Einzelraum, 10–15 Möbelstücke, funktionales Layout",
         "Nur kommerziell sichere Werkzeuge – keine AGPL, keine Umsatzgrenzen",
         "Saubere, wiederholbare Meshes über einen kuratierten Katalog"],
        ["Objective: image→IFC pipeline for office furniture (not full-building BIM)",
         "Scope: single room, 10–15 items, functional layout",
         "Only commercially-safe tools — no AGPL, no revenue caps",
         "Clean, repeatable meshes via a curated catalog"], G, nxt())

    slide_content(prs, "Das Problem — Einzelbild-Grenzen", "The Problem — Single-view Limits",
        ["Einzelbild-Rekonstruktion ist mathematisch unterbestimmt",
         "1) Asymmetrische Beine – kein Symmetrie-Prior",
         "2) Halluzinierte Rück-/Unterseite – kein Bildinhalt dort",
         "3) Nicht-deterministisch – gleiches Foto, andere Meshes",
         "4) Verrauschte Topologie – Löcher, nicht-mannigfaltig",
         "Mesh ≠ semantisches IFC"],
        ["Single-view reconstruction is mathematically ill-posed",
         "1) Asymmetric legs — no symmetry prior",
         "2) Hallucinated back/underside — no image info there",
         "3) Non-deterministic — same photo, different meshes",
         "4) Noisy topology — holes, non-manifold",
         "Mesh ≠ semantic IFC"], G, nxt(), size=13)

    slide_content(prs, "Energie- & Kostenrechnung", "Energy & Cost Calculation",
        ["Strom (Baden-Württemberg 2026): €0,25/kWh",
         "GPU ~0,45–0,6 kW × 24/7 ≈ 324–432 kWh/Monat → €80–108 Strom",
         "Hetzner GEX44: €184/Monat · On-Prem Heilbronn: ~€175/Monat",
         "Pro Raum: ~€0,019 bei 10.000 Räumen/Monat",
         "Lizenzgebühren: €0 – für immer (MIT/Apache/CC-BY)"],
        ["Electricity (Baden-Württemberg 2026): €0.25/kWh",
         "GPU ~0.45–0.6 kW × 24/7 ≈ 324–432 kWh/month → €80–108 power",
         "Hetzner GEX44: €184/month · on-prem Heilbronn: ~€175/month",
         "Per room: ~€0.019 at 10,000 rooms/month",
         "Licence royalties: €0 — forever (MIT/Apache/CC-BY)"], G, nxt(), size=13)

    slide_content(prs, "KI-Ansätze auf einen Blick", "AI Approaches at a Glance",
        ["Retrieval-first: nur bei Katalog-Miss generieren",
         "Retrieval: DINOv2 / SigLIP 2 (Apache-2.0)",
         "Erkennung + Segmentierung: Grounding DINO + SAM 2.1",
         "Tiefe / Maßstab: Depth Anything V2 Small",
         "Generative Reserve: SAM 3D / Stable Fast 3D / TRELLIS / TripoSR"],
        ["Retrieval-first: generate only on a catalog miss",
         "Retrieval: DINOv2 / SigLIP 2 (Apache-2.0)",
         "Detection + segmentation: Grounding DINO + SAM 2.1",
         "Depth / scale: Depth Anything V2 Small",
         "Generative fallback: SAM 3D / Stable Fast 3D / TRELLIS / TripoSR"], G, nxt(), size=13)

    slide_content(prs, "HuggingFace & Lizenzen", "HuggingFace & Licences",
        ["Sichere HF-Modelle: Apache-2.0 / MIT / SAM-Lizenz",
         "Fallen vermeiden:",
         "  Hunyuan3D-2 – in der EU ausgeschlossen",
         "  YOLOv8 – AGPL (Copyleft)",
         "  Depth Anything Base/Large – CC-BY-NC",
         "  Stable Fast 3D – Grenze bei $1 Mio. Umsatz",
         "Prinzip: SCS verkauft Ergebnisse, nicht Gewichte"],
        ["Safe HF models: Apache-2.0 / MIT / SAM license",
         "Traps to avoid:",
         "  Hunyuan3D-2 — excluded in the EU",
         "  YOLOv8 — AGPL (copyleft)",
         "  Depth Anything Base/Large — CC-BY-NC",
         "  Stable Fast 3D — $1M revenue cap",
         "Principle: SCS sells outputs, not weights"], G, nxt(), size=13)

    slide_content(prs, "Der Strategiewechsel — Retrieval + Layout", "The Pivot — Retrieval + Layout",
        ["Statt halluzinieren: Foto → nächstes sauberes Katalog-Mesh",
         "DINOv2 + FAISS über 400 echte ABO-Produktmodelle",
         "Deterministisch, professionell, wiederholbar",
         "Generierung nur als Reserve für Nicht-Katalog-Objekte"],
        ["Instead of hallucinating: photo → nearest clean catalog mesh",
         "DINOv2 + FAISS over 400 real ABO product models",
         "Deterministic, professional, repeatable",
         "Generation only as a fallback for non-catalog objects"], G, nxt())

    # ---- PART 2 — Dimitrius ----
    slide_divider(prs, "Teil 2 — System & Ergebnisse", "Part 2 — System & Results", "Dimitrius")

    slide_figure(prs, "Systemüberblick", "System Overview", "results_plate.png",
        "Gesamtüberblick: Layouts, Kapazitätsgrenze und Foto→3D-Genauigkeit",
        "Overview: layouts, capacity boundary and photo→3D accuracy", D, nxt())

    slide_content(prs, "Der 400-Objekt-Katalog", "The 400-item Catalog",
        ["Amazon Berkeley Objects (ABO): 400 echte Produktmodelle",
         "8 Kategorien × 50, Lizenz CC-BY-4.0",
         "Echte metrische Maße (H×B×T in Metern)",
         "Retrieval: DINOv2 + FAISS",
         "Einzelauswahl mit farbigen Vorschaubildern"],
        ["Amazon Berkeley Objects (ABO): 400 real product models",
         "8 categories × 50, licence CC-BY-4.0",
         "Real metric dimensions (H×W×D in metres)",
         "Retrieval: DINOv2 + FAISS",
         "Per-item picker with colored previews"], D, nxt())

    slide_content(prs, "Layout-Engine — 3 Schichten", "Layout Engine — 3 Layers",
        ["Schicht 1 – Regelpakete (Neufert 6 m²/Arbeitsplatz, ADA 0,915 m, Tür 0,90 m)",
         "Schicht 2 – CP-SAT-Packung, 10-cm-Raster, Wand-Affinität",
         "Schicht 3 – funktionale Verankerung + Sitzausrichtung",
         "Stuhl→Schreibtisch, Monitor/Lampe→Tisch; Stuhl blickt zum Tisch",
         "Auf den Zentimeter verifiziert"],
        ["Layer 1 — rule packs (Neufert 6 m²/workstation, ADA 0.915 m, door 0.90 m)",
         "Layer 2 — CP-SAT packing, 10 cm grid, wall-affinity",
         "Layer 3 — functional anchoring + seat-facing",
         "Chair→desk, monitor/lamp→desk; chair faces the desk",
         "Verified to the centimetre"], D, nxt(), size=13)

    slide_figure(prs, "Layout in Aktion", "Layout in Action", "fig02_office_team_montage.png",
        "Drei-Arbeitsplatz-Büro: Stühle zum Schreibtisch, Lager an den Wänden, Mitte frei",
        "Three-workstation office: chairs face desks, storage on walls, centre open", D, nxt())

    slide_figure(prs, "Randbedingungen & Barrierefreiheit", "Constraints & Accessibility",
        "fig03_office_obstacles_montage.png",
        "Säule + Türfreiraum (links) · ADA breitere Gänge (rechts)",
        "Column + door keep-clear (left) · ADA wider aisles (right)", D, nxt(),
        img2="fig04_office_ada_montage.png")

    slide_figure(prs, "Verallgemeinerung", "Generalization", "fig05_living_room_montage.png",
        "Wohnzimmer (links) · dichter Arbeitsraum (rechts) – Regelpakete pro Raumtyp",
        "Living room (left) · dense workspace (right) — per-room-type rule packs", D, nxt(),
        img2="fig06_workspace_dense_montage.png")

    slide_figure(prs, "Kapazitätsgrenze", "Capacity Boundary", "fig08_capacity_sweep.png",
        "Raumgröße × Arbeitsplätze – 4×3→2, 5×4→3, 6×5→4, 8×6→6; skaliert mit der Fläche",
        "Room size × workstations — 4×3→2, 5×4→3, 6×5→4, 8×6→6; scales with area", D, nxt())

    slide_content(prs, "Die Web-App", "The Web App",
        ["Flask-Backend + xeokit-3D-Betrachter (WebGL)",
         "Ablauf: Auswählen → Generieren → Vorschau → Export",
         "Export: CSV / GLB / IFC4",
         "Flüchtig: nichts wird gespeichert, bis exportiert wird"],
        ["Flask backend + xeokit 3D viewer (WebGL)",
         "Flow: pick → Generate → preview → Export",
         "Export: CSV / GLB / IFC4",
         "Ephemeral: nothing is saved until you export"], D, nxt())

    slide_content(prs, "Genauigkeitsbewertung — Methode", "Accuracy Evaluation — Method",
        ["ABO-Meshes als Ground Truth (wir besitzen sie)",
         "Foto rendern → rekonstruieren → mit Original vergleichen",
         "Metriken: Chamfer-Distanz + F-Score (τ=0,02)",
         "Multi-Seed-ICP-Ausrichtung, deterministisch (seed)",
         "Kalibrierung: Identität F=1,0 · anderes Objekt F=0,18"],
        ["ABO meshes as ground truth (we own them)",
         "Render photo → reconstruct → compare to original",
         "Metrics: Chamfer distance + F-score (τ=0.02)",
         "Multi-seed ICP alignment, deterministic (seeded)",
         "Calibration: identity F=1.0 · different object F=0.18"], D, nxt(), size=13)

    slide_figure(prs, "Genauigkeit — Ergebnis", "Accuracy — Result", "fig09_accuracy_triposr.png",
        "TripoSR vs. ABO: Präzision ~0,81 ≫ Recall ~0,09 – die Einzelbild-Grenze (Chamfer 0,169, F 0,155)",
        "TripoSR vs ABO: precision ~0.81 ≫ recall ~0.09 — the single-view ceiling (chamfer 0.169, F 0.155)",
        D, nxt())

    slide_content(prs, "4-Wege-Vergleich (nächster Schritt)", "4-way Bake-off (next)",
        ["Gleiche Metrik bewertet alle vier Generatoren",
         "TripoSR · InstantMesh · TRELLIS · SAM 3D",
         "Auf RunPod-GPU (A40), ~$10 gesamt",
         "Ergebnis: Vergleichstabelle Modell × Genauigkeit"],
        ["Same metric scores all four generators",
         "TripoSR · InstantMesh · TRELLIS · SAM 3D",
         "On a RunPod GPU (A40), ~$10 total",
         "Output: comparison table model × accuracy"], D, nxt())

    slide_content(prs, "Ergebnisse", "Outcomes",
        ["Funktionierendes Retrieval-+-Layout-System",
         "Einzelbild-Grenze gemessen, nicht nur behauptet",
         "€0 Lizenzgebühren · ~€185/Monat · <€0,02/Raum",
         "Reproduzierbare Abbildungen + IFC/CSV/GLB-Export"],
        ["Working retrieval-+-layout system",
         "Single-view limit measured, not just asserted",
         "€0 royalties · ~€185/month · <€0.02/room",
         "Reproducible figures + IFC/CSV/GLB export"], D, nxt())

    slide_content(prs, "Ausblick", "Future Work",
        ["Multi-View + Maßstabskalibrierung (über die Grenze hinaus)",
         "Foto→Retrieval direkt in der App",
         "IFC4-Validierung (Eigenschaftssätze)",
         "Desktop-App (PySide, .exe)"],
        ["Multi-view + scale calibration (past the ceiling)",
         "Photo→retrieval directly in the app",
         "IFC4 validation (property sets)",
         "Desktop app (PySide, .exe)"], D, nxt())

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"-> {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, {OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
